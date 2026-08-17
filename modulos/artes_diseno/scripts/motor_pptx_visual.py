#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Inyecta texto en plantillas PPTX a partir de un JSON de etiquetas y exporta a PDF.

Uso:
    motor_pptx_visual.py --plantilla base.pptx --input datos_diseno.json --salida diseno_generado.pptx
    motor_pptx_visual.py --plantilla base.pptx --input datos.json --salida salida.pptx --no-pdf
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.error
import urllib.request

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

LLAVE_RE = re.compile(r"\{\{[A-Za-z0-9_ÁÉÍÓÚÑáéíóúñ ]+?\}\}")
USER_AGENT = "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0 Safari/537.36"

ADVERTENCIAS = []


def descargar_imagen(url, timeout=30):
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        data = resp.read()
        ctype = resp.headers.get("Content-Type", "")
    ext = "jpg"
    if "png" in ctype:
        ext = "png"
    elif "gif" in ctype:
        ext = "gif"
    elif "webp" in ctype:
        ext = "webp"
    return data, ext


def etiquetas_de_picture(shape):
    etiquetas = []
    nombre = (shape.name or "").strip()
    if LLAVE_RE.search(nombre):
        etiquetas.extend(LLAVE_RE.findall(nombre))
    cNvPr = shape._element.find(".//" + qn("p:cNvPr"))
    if cNvPr is not None:
        descr = (cNvPr.get("descr") or "").strip()
        if LLAVE_RE.search(descr):
            etiquetas.extend(LLAVE_RE.findall(descr))
    return etiquetas


def reemplazar_imagen(picture, data, ext):
    with tempfile.NamedTemporaryFile(suffix="." + ext, delete=False) as f:
        f.write(data)
        tmp_path = f.name
    try:
        image_part, rId = picture.part.get_or_add_image_part(tmp_path)
    finally:
        os.unlink(tmp_path)
    blips = picture._element.findall(".//" + qn("a:blip"))
    if not blips:
        raise ValueError("La imagen no tiene un blip asociado.")
    blips[0].set(qn("r:embed"), rId)


def procesar_picture(shape, mapeo, encontradas):
    for etiqueta in etiquetas_de_picture(shape):
        valor = mapeo.get(etiqueta)
        if not valor:
            continue
        if isinstance(valor, str) and valor.startswith(("http://", "https://")):
            try:
                data, ext = descargar_imagen(valor)
                reemplazar_imagen(shape, data, ext)
                encontradas.add(etiqueta)
                print(f"Imagen reemplazada: {etiqueta} <- {valor}")
            except (urllib.error.URLError, ValueError, OSError) as e:
                ADVERTENCIAS.append(f"No se pudo descargar {etiqueta} ({valor}): {e}")
        else:
            ADVERTENCIAS.append(f"El valor de {etiqueta} no es una URL válida (http/https).")


def reemplazar_parrafo(parrafo, mapeo, encontradas):
    texto = "".join(run.text for run in parrafo.runs)
    if not texto:
        return
    nuevo = texto
    for clave, valor in mapeo.items():
        if clave in nuevo:
            nuevo = nuevo.replace(clave, valor)
            encontradas.add(clave)
    if nuevo != texto:
        runs = parrafo.runs
        if runs:
            runs[0].text = nuevo
            for run in runs[1:]:
                run.text = ""


def reemplazar_text_frame(tf, mapeo, encontradas):
    for parrafo in tf.paragraphs:
        reemplazar_parrafo(parrafo, mapeo, encontradas)


def reemplazar_tabla(tabla, mapeo, encontradas):
    for fila in tabla.rows:
        for celda in fila.cells:
            reemplazar_text_frame(celda.text_frame, mapeo, encontradas)


def reemplazar_shapes(shapes, mapeo, encontradas):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            reemplazar_shapes(shape.shapes, mapeo, encontradas)
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            procesar_picture(shape, mapeo, encontradas)
            continue
        if shape.has_text_frame:
            reemplazar_text_frame(shape.text_frame, mapeo, encontradas)
        if shape.has_table:
            reemplazar_tabla(shape.table, mapeo, encontradas)


def etiquetas_pendientes(presentacion):
    pendientes = set()
    for slide in presentacion.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texto = "".join(r.text for r in p.runs)
                    pendientes.update(LLAVE_RE.findall(texto))
    return pendientes


def main():
    ap = argparse.ArgumentParser(description="Reemplaza etiquetas {{LLAVE}} en una plantilla PPTX.")
    ap.add_argument("--plantilla", required=True, help="Ruta a la plantilla .pptx base")
    ap.add_argument("--input", default="datos_diseno.json", help="Ruta al JSON de etiquetas")
    ap.add_argument("--salida", default="diseno_generado.pptx", help="Ruta del .pptx de salida")
    ap.add_argument("--no-pdf", action="store_true", help="No convertir el resultado a PDF")
    args = ap.parse_args()

    with open(args.input, encoding="utf-8") as f:
        mapeo = json.load(f)
    if not isinstance(mapeo, dict):
        sys.exit("El JSON debe ser un objeto con pares etiqueta: texto.")

    presentacion = Presentation(args.plantilla)
    encontradas = set()
    for slide in presentacion.slides:
        reemplazar_shapes(slide.shapes, mapeo, encontradas)
        if slide.has_notes_slide:
            reemplazar_text_frame(slide.notes_slide.notes_text_frame, mapeo, encontradas)

    if not args.salida.lower().endswith(".pptx"):
        args.salida += ".pptx"
    os.makedirs(os.path.dirname(os.path.abspath(args.salida)), exist_ok=True)
    presentacion.save(args.salida)
    print(f"PPTX generado: {args.salida}")

    no_usadas = sorted(k for k in mapeo if k not in encontradas)
    if no_usadas:
        print(f"Advertencia: etiquetas del JSON no encontradas en la plantilla: {', '.join(no_usadas)}")
    for advertencia in ADVERTENCIAS:
        print(f"Advertencia: {advertencia}")
    restantes = etiquetas_pendientes(presentacion)
    if restantes:
        print(f"Advertencia: etiquetas sin reemplazo en la plantilla: {', '.join(sorted(restantes))}")

    if not args.no_pdf:
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            sys.exit("LibreOffice no encontrado; no se pudo generar el PDF.")
        out_dir = os.path.dirname(os.path.abspath(args.salida))
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, args.salida],
            check=True,
        )
        base = os.path.splitext(os.path.basename(args.salida))[0]
        print(f"PDF generado: {os.path.join(out_dir, base + '.pdf')}")


if __name__ == "__main__":
    main()